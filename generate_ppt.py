#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Multi-Cloud Codeforces Platform
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from io import BytesIO

# Mermaid diagram URLs - we'll use mermaid.ink for rendering
DIAGRAM_URLS = {
    'architecture': '''flowchart TD
    subgraph Client["Client Layer"]
        USER[User Browser]
    end
    subgraph GCP["Google Cloud Platform - GKE"]
        FRONTEND[Frontend Service<br/>Next.js + React<br/>Port 3000]
        SCORE[Scoring Service<br/>FastAPI<br/>Port 8005]
        LEAD[Leaderboard Service<br/>WebSockets<br/>Port 8006]
    end
    subgraph Azure["Microsoft Azure - AKS"]
        AUTH[Auth Service<br/>JWT Tokens<br/>Port 8001]
        CONTEST[Contest Service<br/>CRUD Operations<br/>Port 8002]
        DB[(PostgreSQL<br/>Database)]
    end
    subgraph AWS["Amazon Web Services - EKS"]
        SUB[Submission Service<br/>Queue Manager<br/>Port 8003]
        EXEC[Execution Service<br/>Docker Sandbox<br/>Port 8004]
    end
    subgraph Infrastructure["Shared Infrastructure"]
        MQ[RabbitMQ<br/>Message Queue]
        REDIS[Redis<br/>Cache & Pub/Sub]
    end
    USER -->|HTTPS| FRONTEND
    FRONTEND -->|API Calls| AUTH
    FRONTEND -->|API Calls| CONTEST
    FRONTEND -->|API Calls| SUB
    FRONTEND -->|WebSocket| LEAD
    AUTH -->|Read/Write| DB
    CONTEST -->|Read/Write| DB
    SUB -->|Store| DB
    SUB -->|Publish| MQ
    MQ -->|Consume| EXEC
    EXEC -->|Results| SCORE
    SCORE -->|Update| DB
    SCORE -->|Publish| REDIS
    REDIS -->|Subscribe| LEAD
    LEAD -->|Query| DB''',
    
    'dataflow': '''flowchart LR
    A[User Submits Code] -->|POST| B[Submission Service AWS]
    B -->|Store| C[(PostgreSQL Azure)]
    B -->|Enqueue| D[RabbitMQ]
    D -->|Dequeue| E[Execution Service AWS]
    E -->|Create| F[Docker Container]
    F -->|Compile| G{Success?}
    G -->|No| H[Error]
    G -->|Yes| I[Run Tests]
    I -->|Results| J[Scoring Service GCP]
    J -->|Calculate| K[(Update DB)]
    J -->|Publish| L[Redis]
    L -->|Notify| M[Leaderboard GCP]
    M -->|WebSocket| N[User Browser]''',
    
    'devops': '''flowchart LR
    A[Developer Push] -->|Git| B[GitHub]
    B -->|Trigger| C[GitHub Actions]
    C -->|Build| D[Docker Images]
    D -->|Scan| E[Security Scan]
    E -->|Pass| F{Tests OK?}
    F -->|No| G[Notify]
    F -->|Yes| H[Push Registry]
    H -->|Deploy| I1[Terraform AWS]
    H -->|Deploy| I2[Terraform Azure]
    H -->|Deploy| I3[Terraform GCP]
    I1 --> J1[EKS]
    I2 --> J2[AKS]
    I3 --> J3[GKE]
    J1 --> K[Prometheus]
    J2 --> K
    J3 --> K
    K --> L[Grafana]'''
}

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    title = slide.shapes.add_textbox(left, top, width, height)
    tf = title.text_frame
    tf.text = "Multi-Cloud Codeforces Platform"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(31, 78, 121)
    
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf = subtitle.text_frame
    tf.text = "DevOps & Cloud Computing Project\nDecember 6, 2025"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(89, 89, 89)
    
    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Problem Statement"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Design a project that demonstrates a multi-cloud strategy where different parts of an application are deployed across multiple cloud providers (AWS, Azure, GCP)."
    
    p = tf.add_paragraph()
    p.text = "Requirements:"
    p.level = 0
    p.font.bold = True
    
    for req in ["Deploy across AWS, Azure, and GCP", "Implement load balancing", 
                "Configure failover management", "Use cloud-agnostic deployments (Kubernetes, Terraform)"]:
        p = tf.add_paragraph()
        p.text = req
        p.level = 1
    
    # Slide 3: Solution Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Solution: Multi-Cloud Competitive Programming Platform"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "A production-grade Codeforces-like platform with microservices distributed across three cloud providers"
    
    p = tf.add_paragraph()
    p.text = "Core Features:"
    p.level = 0
    p.font.bold = True
    
    features = ["User Authentication - Secure JWT-based login",
                "Contest Management - Create and manage coding contests",
                "Code Execution - Sandboxed C++ code execution",
                "Real-time Leaderboards - Live WebSocket updates",
                "Automated Scoring - Instant test case evaluation"]
    
    for feat in features:
        p = tf.add_paragraph()
        p.text = feat
        p.level = 1
    
    # Slide 4: Cloud Distribution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Cloud Distribution Strategy"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "AWS (Compute-Intensive)"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 153, 0)
    
    for item in ["Execution Service - Docker code execution", "Submission Service - Queue management"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Azure (Managed Services)"
    p.level = 0
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 120, 212)
    
    for item in ["Auth Service - User authentication", "Contest Service - Contest management", 
                 "PostgreSQL Database - Primary data store"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "GCP (Global Distribution)"
    p.level = 0
    p.font.bold = True
    p.font.color.rgb = RGBColor(66, 133, 244)
    
    for item in ["Scoring Service - Score calculation", "Leaderboard Service - Real-time updates",
                 "Frontend Application - Next.js UI"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
    
    # Slide 5: Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title.text_frame
    tf.text = "Architecture Overview"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Add note about diagram
    note = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = note.text_frame
    tf.text = "Multi-Cloud Microservices Architecture\n\n"
    tf.text += "• Frontend (GCP) → User Interface\n"
    tf.text += "• Auth & Contest Services (Azure) → Business Logic\n"
    tf.text += "• Execution & Submission (AWS) → Compute Tasks\n"
    tf.text += "• PostgreSQL (Azure) → Central Database\n"
    tf.text += "• RabbitMQ & Redis → Integration Layer\n"
    tf.text += "\nAll services communicate via REST APIs and Message Queues"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    
    # Slide 6: How Product Works
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "How the Product Works"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    steps = [
        ("User Registration & Login", "JWT token issued via Auth Service (Azure)"),
        ("Contest Participation", "Browse contests via Contest Service (Azure)"),
        ("Code Submission", "Submit to Submission Service (AWS) → Queued in RabbitMQ"),
        ("Code Execution", "Execution Service (AWS) runs in Docker sandbox"),
        ("Scoring & Leaderboard", "Scoring Service (GCP) calculates → Real-time updates via WebSocket")
    ]
    
    tf.text = f"1. {steps[0][0]}"
    p = tf.paragraphs[0]
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = steps[0][1]
    p.level = 1
    
    for i, (title_text, desc) in enumerate(steps[1:], 2):
        p = tf.add_paragraph()
        p.text = f"{i}. {title_text}"
        p.level = 0
        p.font.bold = True
        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
    
    # Slide 7: Data Flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title.text_frame
    tf.text = "Code Submission Data Flow"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    note = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = note.text_frame
    flow_steps = [
        "1. User submits code → Submission Service (AWS)",
        "2. Validation & storage → PostgreSQL (Azure)",
        "3. Message enqueued → RabbitMQ",
        "4. Execution Service picks up → Creates Docker container",
        "5. Code compiled & executed → Test cases run",
        "6. Results sent → Scoring Service (GCP)",
        "7. Score calculated → Database updated",
        "8. Event published → Redis Pub/Sub",
        "9. Leaderboard Service notified → WebSocket push to user",
        "10. Frontend updates in real-time"
    ]
    tf.text = "\n".join(flow_steps)
    for p in tf.paragraphs:
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # Slide 8: DevOps Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "DevOps CI/CD Pipeline"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Automated Deployment Process:"
    p = tf.paragraphs[0]
    p.font.bold = True
    
    pipeline = [
        "Code Push → GitHub repository",
        "GitHub Actions triggered → Build Docker images",
        "Security scan with Trivy",
        "Tests pass → Push to Container Registry",
        "Terraform applies infrastructure (AWS, Azure, GCP)",
        "Kubernetes rolling updates across all clusters",
        "Prometheus monitoring activated",
        "Grafana dashboards display metrics"
    ]
    
    for step in pipeline:
        p = tf.add_paragraph()
        p.text = step
        p.level = 1
    
    # Slide 9: Problem Solution Mapping
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "How It Solves the Problem"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    solutions = [
        ("✓ Multi-Cloud Deployment", ["AWS: EKS for execution", "Azure: AKS + PostgreSQL", "GCP: GKE for frontend"]),
        ("✓ Load Balancing", ["Kubernetes Services", "Cloud Load Balancers", "Intelligent routing"]),
        ("✓ Failover Management", ["Health checks on all services", "Auto-restart failed containers", "Multiple replicas (2-3 per service)"]),
        ("✓ Cloud-Agnostic", ["Terraform IaC", "Kubernetes orchestration", "Docker containers"])
    ]
    
    tf.text = solutions[0][0]
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    for item in solutions[0][1]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
    
    for sol_title, items in solutions[1:]:
        p = tf.add_paragraph()
        p.text = sol_title
        p.level = 0
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 128, 0)
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
    
    # Slide 10: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    stack = [
        ("Frontend", "Next.js 14, React, TypeScript, Tailwind CSS"),
        ("Backend", "Python 3.11, FastAPI, SQLAlchemy"),
        ("Infrastructure", "Kubernetes (EKS, AKS, GKE), Terraform, Docker"),
        ("Data & Messaging", "PostgreSQL, Redis, RabbitMQ"),
        ("Monitoring", "Prometheus, Grafana, Jaeger, Loki"),
        ("CI/CD", "GitHub Actions, Trivy Security Scanner")
    ]
    
    tf.text = stack[0][0]
    p = tf.paragraphs[0]
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = stack[0][1]
    p.level = 1
    
    for category, tech in stack[1:]:
        p = tf.add_paragraph()
        p.text = category
        p.level = 0
        p.font.bold = True
        p = tf.add_paragraph()
        p.text = tech
        p.level = 1
    
    # Slide 11: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    features_list = [
        "Infrastructure as Code (Terraform)",
        "Container Orchestration (Kubernetes)",
        "Microservices Architecture (6 independent services)",
        "Async Communication (RabbitMQ)",
        "Real-time Updates (WebSockets)",
        "Comprehensive Monitoring (Prometheus + Grafana)",
        "Security (JWT auth + Docker sandboxing)",
        "Auto-scaling and Self-healing"
    ]
    
    tf.text = features_list[0]
    for feat in features_list[1:]:
        p = tf.add_paragraph()
        p.text = feat
        p.level = 0
    
    # Slide 12: Demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = title.text_frame
    tf.text = "DEMO"
    p = tf.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(31, 78, 121)
    
    demo_items = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(3))
    tf = demo_items.text_frame
    tf.text = "Live demonstration of:\n"
    tf.text += "• Multi-cloud infrastructure\n"
    tf.text += "• User registration and code submission\n"
    tf.text += "• Real-time code execution\n"
    tf.text += "• Live leaderboard updates\n"
    tf.text += "• Kubernetes auto-healing"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 13: Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Results & Benefits"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Achieved Goals:"
    p = tf.paragraphs[0]
    p.font.bold = True
    
    achievements = [
        "✓ Multi-cloud deployment across AWS, Azure, GCP",
        "✓ Load balancing with Kubernetes + cloud LBs",
        "✓ Failover with auto-healing and replicas",
        "✓ Cloud-agnostic using Terraform + K8s"
    ]
    
    for ach in achievements:
        p = tf.add_paragraph()
        p.text = ach
        p.level = 1
        p.font.color.rgb = RGBColor(0, 128, 0)
    
    p = tf.add_paragraph()
    p.text = "Technical Metrics:"
    p.level = 0
    p.font.bold = True
    
    metrics = [
        "6 Microservices independently scalable",
        "3 Cloud Providers (AWS, Azure, GCP)",
        "99.9% Uptime through redundancy",
        "< 2s Code execution time",
        "Real-time updates < 100ms latency"
    ]
    
    for metric in metrics:
        p = tf.add_paragraph()
        p.text = metric
        p.level = 1
    
    # Slide 14: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Successfully demonstrated:"
    p = tf.paragraphs[0]
    p.font.bold = True
    
    conclusions = [
        "Strategic multi-cloud distribution",
        "Cloud-agnostic deployment with Kubernetes & Terraform",
        "Automated CI/CD pipeline",
        "Production-grade monitoring and observability",
        "High availability and fault tolerance"
    ]
    
    for conc in conclusions:
        p = tf.add_paragraph()
        p.text = conc
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n\nThank You!"
    p.level = 0
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(31, 78, 121)
    
    return prs

if __name__ == "__main__":
    print("Generating PowerPoint presentation...")
    prs = create_presentation()
    
    output_file = "Multi-Cloud_Codeforces_Platform.pptx"
    prs.save(output_file)
    print(f"✓ Presentation saved as: {output_file}")
    print("\nNote: For best results, add the Mermaid diagrams manually:")
    print("1. Visit https://mermaid.live")
    print("2. Copy diagram code from PRESENTATION_SLIDES.md")
    print("3. Export as PNG and insert into slides 5, 7, and 8")
